"""Generate TelecomGPT AI Agent Architecture PowerPoint (updated)."""

from __future__ import annotations

import sys
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

_OUTPUT_DIR = ROOT / "data" / "reports"

# (title, bullets) — role/responsibility slides use "Agent | Role | Responsibility" pattern
SLIDES: list[tuple[str, list[str]]] = [
    (
        "TelecomGPT — AI Agent Architecture",
        [
            "Domain-specific multi-agent copilot for 5G/LTE RF & Test Engineering",
            "Stack: Adaptive RAG + LangGraph + FastAPI + Next.js (+ Streamlit analytics)",
            "Production: Render 2GB API + Vercel UI",
            "Audience: Senior Test Engineer — logs, KPIs, 3GPP validation, reports",
        ],
    ),
    (
        "Core Concepts",
        [
            "Orchestrator — LangGraph state machine routes each query to specialist agents",
            "Adaptive RAG — retrieval adapts by query: BM25 + vector + live fetch + web",
            "Structured KB — deterministic bands, devices, calculators (fast, traceable)",
            "Session memory — uploads and chat context bound to session_id",
            "Guardrails — input/output filtering, PII redaction, tool allowlists",
            "Artifacts — text, Plotly charts, attach/UE-cap reports, PDF/Excel, .pptx",
        ],
    ),
    (
        "Deployment Architecture",
        [
            "Next.js (Vercel) — primary chat UI, suggestion chips, agent trace",
            "FastAPI (Render 2GB) — POST /ask wraps LangGraph orchestrator",
            "Streamlit (optional) — CSV/log analytics dashboard, not main chat path",
            "Build: ingest_rag.py at deploy → chunks.json (~2,230 BM25 chunks)",
            "Env: TELECOMGPT_VECTOR=1, LIVE_FETCH=1, LOW_MEMORY=0, 8 parallel agents",
        ],
    ),
    (
        "LangGraph Pipeline — Roles",
        [
            "load_memory — assembles session + semantic/episodic/procedural context",
            "guardrails_pre — blocks/redacts unsafe input before planning",
            "plan — keyword (+ optional LLM) planner selects agent list",
            "confidence_gate — clarifies vague queries or proceeds with score",
            "parallel_batch — runs task/retrieval/autonomous agents concurrently",
            "sequential_tail — presentation → synthesizer → verifier",
            "guardrails_post — filters LLM output",
            "save_memory — persists Q&A and successful plans",
        ],
    ),
    (
        "Two User Interaction Paths",
        [
            "Path A — Conversational: chips / chat → POST /ask → LangGraph + RAG",
            "Path B — Deterministic tools: Attach Report / UE Cap Report → dedicated APIs",
            "Path A — best for explain, troubleshoot, open-ended analysis",
            "Path B — best for repeatable checklist scoring + PDF/Excel export",
            "Upload (📎) — POST /api/upload binds CSV/log to session for Path A agents",
        ],
    ),
    (
        "Adaptive Hybrid RAG — Concept",
        [
            "Simple path — glossary, band, device, ARFCN → structured KB (no heavy RAG)",
            "Complex path — explain, spec, fault → hybrid_retrieve + multi-agent",
            "BM25 — static chunks.json (ShareTechnote, RF handbook, sqimway, 3GPP)",
            "Vector — ChromaDB long-term reference + conversation memory",
            "Live fetch — ShareTechnote topic + sqimway band rows + 3GPP TS series",
            "Optional Tavily web — telecom domain bias (sharetechnote, sqimway, 3gpp.org)",
        ],
    ),
    (
        "Memory Architecture — Roles",
        [
            "Short-term (session) — last ~100 chat turns, upload file paths",
            "Semantic — user facts (bands of interest, devices)",
            "Episodic — past Q&A summaries",
            "Procedural — last successful agent plans",
            "Reference (vector) — RAG chunks indexed via POST /api/memory/ingest-rag",
            "User profile — bands, devices, notes per session_id",
        ],
    ),
    (
        "Agent Taxonomy Overview",
        [
            "Task agents — bounded workflows with deterministic tools",
            "Retrieval agents — search structured + unstructured knowledge",
            "Autonomous agents — dynamic tool selection (KB, ReAct, CrewAI, AutoGen)",
            "Orchestration agents — synthesizer merges; verifier cross-checks KB",
            "Total: 22 named agents — GET /api/agents/taxonomy",
        ],
    ),
    (
        "Task Agents — Roles & Responsibilities (1/2)",
        [
            "log_debug | Log analysis | Parse QXDM/QCAT logs; attach/UE-cap hints; protocol stack scan",
            "fault_analysis | Troubleshooting | Symptom→cause→checks from fault catalog + stack context",
            "rf_metrics | RF KPIs | Grade drive-test CSV; RSRP/throughput; RF handbook hints",
            "bts_config | gNB config | Scan BTS exports for NR params; cross-check 3GPP limits",
            "feature_validation | Test plans | Built-in 3GPP feature templates + pass criteria",
            "drive_test | Field test | SLA rules, GPS RF maps when lat/lon present",
        ],
    ),
    (
        "Task Agents — Roles & Responsibilities (2/2)",
        [
            "analytics | Data viz | Kaggle CSV summaries, Plotly chart artifacts",
            "prediction | Trends | KPI correlation and trend analysis",
            "presentation | Reports | Generate PowerPoint via python-pptx",
            "comparison | Diff | Device vs device, technology comparisons",
            "compliance | Regulatory | FCC licensed band checks",
            "deploy | Ops | Production health, engine status",
            "eval | QA | KB smoke/regression tests",
        ],
    ),
    (
        "Retrieval & Autonomous Agents",
        [
            "research | Hybrid RAG | BM25 + vector + live ST/sqimway/3GPP + memory recall",
            "spec | 3GPP lookup | TS-focused hybrid search with citations",
            "telecom_kb | Structured KB | Bands, devices, CA/EN-DC, ARFCN/GSCN/throughput",
            "react | ReAct loop | LLM picks 1–2 tools iteratively",
            "crew | CrewAI | Role-based researcher + RF + compliance crew (hybrid engine)",
            "autogen | AutoGen | Multi-turn autonomous tool calling (hybrid engine)",
        ],
    ),
    (
        "Orchestration Agents",
        [
            "synthesizer | Final answer | Merges agent outputs + RAG + LLM; appends Sources",
            "verifier | Quality gate | Cross-checks answer vs KB agent outputs; flags gaps",
            "Planner (plan node) | Routing | Maps query keywords → agent list; optional LLM refine",
            "Confidence gate | UX safety | Low-confidence short queries → clarification prompt",
        ],
    ),
    (
        "Test Engineer Tools (Path B)",
        [
            "NR SA Attach Report — 14-step checklist on log; PDF/Excel export",
            "UE Capability Report — Enquiry/Information/segmentation scan; PDF/Excel",
            "NR band catalog — 91 bands from sqimway (TS 38.104)",
            "Protocol stack reference — C/U-plane, PHY→NAS, layer→spec map",
            "Power class reference — FR1/FR2, HPUE, RedCap PC7",
            "RF Handbook — curated ShareTechnote RF topics + KPI crosswalk",
        ],
    ),
    (
        "UI → Architecture Mapping",
        [
            "Suggestion chips — intent hints → planner → agent selection",
            "Show agent trace — trace:true → plan, steps, confidence, guardrails",
            "📎 Upload — session-scoped data for log_debug / rf_metrics / analytics",
            "📋 Attach Report — POST /api/nr-sa/attach-report (deterministic scanner)",
            "📡 UE Cap Report — POST /api/nr/ue-capability/report",
        ],
    ),
    (
        "Key Tools (ToolRegistry)",
        [
            "lookup_glossary, lookup_device, lookup_bands, lookup_ca_endc, calc_phy",
            "hybrid_search — adaptive RAG entry point for research/spec agents",
            "memory_search — vector recall by session",
            "csv_summary, analyze_log, run_drive_test_rules, list_kaggle_csvs",
            "generate_presentation, web_search (Tavily), live_reference_fetch",
        ],
    ),
    (
        "Knowledge Sources",
        [
            "telecom_master_db.json — NR/LTE bands, glossary, FCC, comparisons",
            "devices/*.json — S23, S24, iPhone, Pixel capability sheets",
            "nr_bands_catalog.json — 91 NR bands (sqimway / TS 38.104)",
            "Reference JSON — attach checklist, UE cap, protocol stack, power class, RF handbook",
            "RAG chunks — ShareTechnote 5G + RF handbook + sqimway + 3GPP overview",
            "Live: sharetechnote.com, sqimway.com/nr_band.php, 3gpp.org dynareport",
        ],
    ),
    (
        "API Endpoints (Agent-facing)",
        [
            "POST /ask — multi-agent chat (async job for slow queries)",
            "POST /api/upload — CSV/log for session agents",
            "POST /api/rag/reindex — rebuild BM25 chunk store",
            "POST /api/memory/ingest-rag — background vector index",
            "GET /api/memory/ingest-rag/status — poll vector ingest progress",
            "GET /api/agents/taxonomy — full agent map",
            "POST /api/nr-sa/attach-report | /api/nr/ue-capability/report",
        ],
    ),
    (
        "Configuration (Production 2GB)",
        [
            "TELECOMGPT_LOW_MEMORY=0 — full agent parallelism (8 workers)",
            "TELECOMGPT_VECTOR=1 — Chroma hybrid retrieval enabled",
            "TELECOMGPT_LIVE_FETCH=1 — live ShareTechnote/sqimway/3GPP supplement",
            "TELECOMGPT_AUTO_REINDEX=1 — vector ingest on boot (background thread)",
            "TELECOMGPT_ENGINE=hybrid — LangGraph master + CrewAI/AutoGen spokes",
            "OPENAI_API_KEY — gpt-4o-mini synthesis; Ollama optional for local dev",
        ],
    ),
    (
        "Orchestrator Flow (Visual)",
        [
            "User → Next.js → POST /ask → FastAPI",
            "load_memory → guardrails → plan → confidence_gate",
            "parallel_batch (up to 8 agents) → synthesizer → verifier",
            "guardrails_post → save_memory → JSON answer + artifacts + sources",
            "Optional: async /api/jobs/{id} for charts, PPT, heavy analytics",
        ],
    ),
    (
        "Trust & Explainability",
        [
            "Deterministic KB/calculators cite 3GPP clauses (ARFCN, throughput)",
            "RAG answers include ShareTechnote / 3GPP / sqimway Sources URLs",
            "Path B reports — rule-based checklist scores (auditable for test labs)",
            "Agent trace — exposes plan and steps for test engineer validation",
            "Verifier — flags when KB mentions bands/devices missing from answer",
        ],
    ),
    (
        "Summary",
        [
            "TelecomGPT = LangGraph orchestrator + 22 agents + adaptive hybrid RAG",
            "Test Engineer focus: logs, KPIs, attach/UE-cap, 3GPP refs, one-click exports",
            "Simple queries → KB fast path; complex → agents + live RAG + synthesizer",
            "Next.js chat + Streamlit analytics; Render 2GB + Vercel",
            "Full docs: docs/ARCHITECTURE.md, docs/ORCHESTRATION.md",
        ],
    ),
]


def generate_deck() -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"ok": False, "error": "pip install python-pptx"}

    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now(timezone.utc).strftime("%Y%m%d")
    filename = f"TelecomGPT_AI_Agent_Architecture_{ts}.pptx"
    out_path = _OUTPUT_DIR / filename

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TelecomGPT AI Agent Architecture"
    try:
        slide.placeholders[1].text = (
            "Adaptive RAG + LangGraph + FastAPI + Next.js\n"
            "Concepts · Roles · Responsibilities\n"
            f"{datetime.now(timezone.utc).strftime('%B %Y')}"
        )
    except (IndexError, KeyError):
        pass

    # Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate([
        "Core concepts & deployment topology",
        "LangGraph pipeline roles",
        "Adaptive hybrid RAG & memory layers",
        "Agent taxonomy — roles & responsibilities",
        "Test Engineer tools & UI mapping",
        "APIs, config & trust model",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)

    for title, bullets in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet[:400]
            p.level = 0
            p.font.size = Pt(17 if "|" in bullet else 18)

    # Agent category matrix slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agent Category Matrix"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    matrix = [
        "TASK — Execute bounded tools on data (logs, CSV, config, reports)",
        "RETRIEVAL — Search KB + RAG + live references (research, spec)",
        "AUTONOMOUS — LLM-driven tool selection (telecom_kb, react, crew, autogen)",
        "ORCHESTRATION — Merge & verify (synthesizer, verifier)",
        "INFRA NODES — memory, guardrails, planner, confidence (LangGraph graph)",
    ]
    for i, line in enumerate(matrix):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    try:
        slide.placeholders[1].text = (
            "github.com/aqil2020in/telecomgpt\n"
            "API: telecomgpt.onrender.com · UI: telecomgpt.vercel.app"
        )
    except (IndexError, KeyError):
        pass

    prs.save(str(out_path))
    return {
        "ok": True,
        "path": str(out_path),
        "filename": filename,
        "slides": len(SLIDES) + 4,
        "download_url": f"/api/reports/{filename}",
    }


if __name__ == "__main__":
    result = generate_deck()
    if result.get("ok"):
        print(f"Generated: {result['path']}")
        print(f"Slides: {result['slides']}")
        print(f"Download: {result['download_url']}")
    else:
        print(f"Error: {result.get('error')}")
