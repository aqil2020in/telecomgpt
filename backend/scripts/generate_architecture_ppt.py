"""Generate the TelecomGPT architecture overview PowerPoint deck."""

from __future__ import annotations

import sys
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

_OUTPUT_DIR = ROOT / "data" / "reports"

# Each entry: (title, bullets)
ARCHITECTURE_SLIDES: list[tuple[str, list[str]]] = [
    (
        "What is TelecomGPT?",
        [
            "Domain-specific AI assistant for cellular / RF engineering",
            "Covers 5G NR, LTE, bands, devices, CA/EN-DC, 3GPP calculations",
            "Production: chat UI (Vercel) + API (Render) + local dev stack",
            "Built as a multi-agent orchestrator with tools, memory, and RAG",
        ],
    ),
    (
        "Build Journey — 8 Phases",
        [
            "Phase 1: Structured knowledge base (JSON + device sheets)",
            "Phase 2: LangGraph keyword router (deterministic handlers)",
            "Phase 3: BM25 RAG over ShareTechnote + 3GPP references",
            "Phase 4: Multi-agent orchestrator + autonomous planning",
            "Phase 5: Tool-use framework (12 callable tools)",
            "Phase 6: Vector + session memory (ChromaDB)",
            "Phase 7: Analytics + Kaggle 5G datasets",
            "Phase 8: PowerPoint report generation + chat UI artifacts",
        ],
    ),
    (
        "High-Level Architecture",
        [
            "Client: Next.js chat UI (frontend/src/pages/index.tsx)",
            "API: FastAPI backend (backend/app.py)",
            "Brain: LangGraph orchestrator (backend/telecom_ai/orchestrator.py)",
            "Knowledge: telecom_master_db.json + devices/*.json",
            "Retrieval: BM25 RAG (2230 chunks) + ChromaDB vector memory",
            "Output: text answers, sources, Plotly charts, .pptx downloads",
        ],
    ),
    (
        "Multi-Agent Orchestrator Flow",
        [
            "START → load_memory → plan → orchestrator (supervisor loop)",
            "Specialists: telecom_kb | research | analytics | presentation",
            "Each agent calls tools; outputs collected in shared state",
            "synthesizer merges agent outputs + LLM + RAG sources",
            "save_memory persists session + vector store → END",
            "Legacy mode: TELECOMGPT_MODE=legacy uses keyword router (graph.py)",
        ],
    ),
    (
        "Specialist Agents",
        [
            "telecom_kb — bands, devices, CA/EN-DC, ARFCN/GSCN/throughput",
            "research — RAG search + vector memory recall",
            "analytics — Kaggle CSV summaries, RF column detection, logs",
            "presentation — builds PowerPoint via python-pptx",
            "synthesizer — final LLM answer with citations and artifacts",
        ],
    ),
    (
        "Autonomous Planning System",
        [
            "File: backend/telecom_ai/planning.py",
            "Rule-based planner detects intent (PPT, analytics, KB, research)",
            "Produces step list: which agents run and in what order",
            "Optional LLM refinement via TELECOMGPT_LLM_PLAN=1",
            "Example: PPT request → research → telecom_kb → presentation → synthesizer",
        ],
    ),
    (
        "Tool-Use Framework",
        [
            "File: backend/telecom_ai/tools.py — ToolRegistry pattern",
            "lookup_glossary, lookup_device, lookup_ca_endc, calc_phy",
            "lookup_bands, rag_search, memory_search",
            "csv_summary, detect_rf_columns, analyze_log, list_kaggle_csvs",
            "generate_presentation — exposed via GET /api/tools",
        ],
    ),
    (
        "Memory System",
        [
            "Session memory: JSON per session (backend/memory/session_memory.py)",
            "Vector memory: ChromaDB + lightweight fallback index",
            "Conversation turns saved after each /ask response",
            "RAG chunks can be indexed: POST /api/memory/ingest-rag",
            "Recalled at start of each orchestrator run (load_memory node)",
        ],
    ),
    (
        "Knowledge Base Layer",
        [
            "telecom_master_db.json — NR/LTE bands, glossary, comparisons, FCC",
            "devices/*.json — S23, S24, iPhone, Pixel band/combo sheets",
            "Calculators (3GPP-grounded):",
            "  • arfcn_calculator.py — NR-ARFCN ↔ frequency",
            "  • gscn_calculator.py — GSCN ↔ SSB frequency",
            "  • throughput_calculator.py — peak NR data rate",
        ],
    ),
    (
        "RAG — Reference Retrieval",
        [
            "Ingest: backend/scripts/ingest_rag.py fetches ShareTechnote pages",
            "Store: backend/data/rag/chunks.json (~2230 chunks, 80+ pages)",
            "Search: BM25 lexical retrieval (backend/rag/retrieve.py)",
            "Merged into LLM context with KB excerpts and sources appended",
            "Reindex: POST /api/rag/reindex (dev/admin)",
        ],
    ),
    (
        "Analytics & Kaggle Data",
        [
            "Streamlit dashboard: analytics/app.py (CSV, logs, Plotly charts)",
            "FastAPI routes: /api/analytics/csv/*, /logs/analyze, /kaggle/*",
            "Curated Kaggle 5G datasets in backend/data/kaggle/",
            "Download script: python scripts/download_kaggle.py download --all",
            "RF auto-detection: RSRP, lat/lon, throughput column aliases",
        ],
    ),
    (
        "PowerPoint Report Generation",
        [
            "Module: backend/ppt/generator.py (python-pptx)",
            "Chat: 'Generate a PowerPoint report on …' triggers presentation agent",
            "API: POST /api/ppt/generate → GET /api/reports/{filename}",
            "Frontend shows green Download button on assistant messages",
            "This deck generated by backend/scripts/generate_architecture_ppt.py",
        ],
    ),
    (
        "Tech Stack & Resources",
        [
            "Backend: Python 3.12, FastAPI, Uvicorn, LangGraph, LangChain-core",
            "LLM: OpenAI GPT-4o-mini and/or Ollama (llama3.1) — TELECOMGPT_LLM=auto",
            "Data: Pandas, Plotly, python-pptx, ChromaDB, Requests",
            "Frontend: Next.js (pages router), TypeScript, React",
            "Deploy: Render (API) + Vercel (UI) — render.yaml blueprint",
            "External refs: ShareTechnote, 3GPP, Kaggle 5G datasets",
        ],
    ),
    (
        "Key Source Files",
        [
            "backend/app.py — REST API entry",
            "backend/telecom_ai/core.py — TelecomAI facade",
            "backend/telecom_ai/orchestrator.py — multi-agent graph",
            "backend/telecom_ai/agents/specialists.py — agent logic",
            "backend/telecom_ai/planning.py — planner",
            "backend/telecom_ai/tools.py — tool registry",
            "docs/ARCHITECTURE.md — Mermaid diagram + full reference",
        ],
    ),
    (
        "API Endpoints",
        [
            "POST /ask — chat with history, session_id, artifacts",
            "GET /api/tools — list agent tools",
            "POST /api/ppt/generate — create .pptx report",
            "GET /api/reports/{file} — download report",
            "POST /api/analytics/csv/summary — CSV stats",
            "POST /api/memory/ingest-rag — index RAG into vector DB",
            "GET /api/health, /api/devices, /api/bands",
        ],
    ),
    (
        "Configuration (Environment Variables)",
        [
            "OPENAI_API_KEY — LLM synthesis (Render env, never commit)",
            "TELECOMGPT_LLM — auto | openai | ollama",
            "TELECOMGPT_MODE — orchestrator (default) | legacy",
            "TELECOMGPT_LLM_PLAN — 1 = LLM plan refinement",
            "RAG_TOP_K — number of RAG chunks retrieved (default 5)",
            "NEXT_PUBLIC_API_URL — frontend → backend URL (Vercel)",
        ],
    ),
    (
        "Local Execution",
        [
            "Backend: cd backend && pip install -r requirements.txt",
            "         uvicorn app:app --port 8000",
            "Frontend: cd frontend && npm install && npm run dev",
            "         NEXT_PUBLIC_API_URL=http://localhost:8000",
            "Analytics UI: streamlit run analytics/app.py",
            "Kaggle: kaggle auth login && python scripts/download_kaggle.py download --all",
            "Swagger docs: http://localhost:8000/docs",
        ],
    ),
    (
        "Production Deployment",
        [
            "GitHub: github.com/aqil2020in/telecomgpt",
            "API: https://telecomgpt.onrender.com (Render Starter)",
            "UI: https://telecomgpt.vercel.app",
            "Render build: pip install -r backend/requirements.txt",
            "Note: Kaggle data gitignored — download locally or at build time",
            "ChromaDB embeddings ~80 MB — consider plan limits on Render Starter",
        ],
    ),
    (
        "Roadmap — What's Next",
        [
            "Analytics charts embedded in chat UI (Plotly from Kaggle CSVs)",
            "File upload in chat — drive-test CSV / UE log analysis",
            "GPS RF coverage maps (Leaflet) from lat/lon columns",
            "QXDM / vendor-specific log parsers",
            "ML prediction lite — RSRP interpolation, throughput forecasting",
            "Production redeploy with orchestrator + python-pptx + chromadb",
        ],
    ),
    (
        "Summary",
        [
            "TelecomGPT = orchestrator + 4 specialist agents + 12 tools",
            "Grounded in JSON KB, BM25 RAG, vector memory, and 3GPP math",
            "Delivers chat answers, source citations, analytics, and PowerPoint reports",
            "Full architecture diagram: docs/ARCHITECTURE.md",
            "Questions / demos: POST /ask with trace=true for agent step visibility",
        ],
    ),
]


def generate_architecture_deck() -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return {
            "ok": False,
            "error": "python-pptx not installed. Run: pip install python-pptx",
        }

    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now(timezone.utc).strftime("%Y%m%d")
    filename = f"TelecomGPT_Architecture_Build_Guide_{ts}.pptx"
    out_path = _OUTPUT_DIR / filename

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TelecomGPT Architecture"
    try:
        slide.placeholders[1].text = (
            "How We Built It — Steps, Resources, Code, Tools & Execution\n"
            f"{datetime.now(timezone.utc).strftime('%B %Y')}"
        )
    except (IndexError, KeyError):
        pass

    # Agenda slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    agenda = [
        "Project overview & build phases",
        "Multi-agent orchestrator architecture",
        "Planning, tools, memory & RAG",
        "Analytics, Kaggle data & PowerPoint generation",
        "Tech stack, APIs & deployment",
        "Local execution & roadmap",
    ]
    for i, item in enumerate(agenda):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)

    # Content slides
    for title, bullets in ARCHITECTURE_SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Handle sub-bullets (lines starting with spaces or •)
            if bullet.startswith("  •") or bullet.startswith("  -"):
                p.text = bullet.strip()
                p.level = 1
                p.font.size = Pt(16)
            else:
                p.text = bullet[:350]
                p.level = 0
                p.font.size = Pt(18)

    # Architecture flow text slide (ASCII-style for PPT)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Orchestrator Pipeline (Visual)"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    flow_lines = [
        "User Query (Next.js Chat UI)",
        "    ↓  POST /ask",
        "load_memory → plan → orchestrator",
        "    ↓ dispatch loop",
        "telecom_kb | research | analytics | presentation",
        "    ↓ tools: KB, RAG, CSV, PPT",
        "synthesizer (LLM + sources + artifacts)",
        "    ↓",
        "save_memory → Response + optional .pptx download",
    ]
    for i, line in enumerate(flow_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        try:
            p.font.name = "Consolas"
        except Exception:
            pass

    # Closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    try:
        slide.placeholders[1].text = (
            "TelecomGPT — github.com/aqil2020in/telecomgpt\n"
            "docs/ARCHITECTURE.md | telecomgpt.vercel.app"
        )
    except (IndexError, KeyError):
        pass

    prs.save(str(out_path))
    slide_count = len(ARCHITECTURE_SLIDES) + 4  # title, agenda, flow, closing

    return {
        "ok": True,
        "path": str(out_path),
        "filename": filename,
        "slides": slide_count,
        "download_url": f"/api/reports/{filename}",
    }


if __name__ == "__main__":
    result = generate_architecture_deck()
    if result.get("ok"):
        print(f"Generated: {result['path']}")
        print(f"Slides: {result['slides']}")
    else:
        print(f"Error: {result.get('error')}")
